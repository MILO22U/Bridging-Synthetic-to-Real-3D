"""Generate PowerPoint presentation for AI 535 project."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

VIZ = './visualizations/presentation'
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
BG_MED = RGBColor(0x16, 0x21, 0x3e)
ACCENT = RGBColor(0x0f, 0x7b, 0xc0)
ACCENT2 = RGBColor(0xe4, 0x3f, 0x5a)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
GOLD = RGBColor(0xFF, 0xD7, 0x00)


def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_slide(slide, left, top, width, height, bullets, font_size=16,
                     color=WHITE, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(8)
        p.level = 0
    return tf


def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        kwargs = {}
        if width: kwargs['width'] = Inches(width)
        if height: kwargs['height'] = Inches(height)
        slide.shapes.add_picture(path, Inches(left), Inches(top), **kwargs)
        return True
    return False


def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                    Inches(width), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


# ================================================================
# SLIDE 1: Title
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 1.5, 1.5, 10, 1.2,
             "Bridging the Synthetic-to-Real Gap\nin 3D Object Reconstruction",
             font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, 4, 3.2, 5.3)

add_text_box(slide, 1.5, 3.5, 10, 0.8,
             "Single-Image 3D Point Cloud Prediction with Domain Adaptation",
             font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1.5, 5.0, 10, 0.5,
             "Mrinal Bharadwaj",
             font_size=22, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1.5, 5.5, 10, 0.5,
             "AI 535 \u2014 Deep Learning  |  March 18, 2026",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 2: Problem & Motivation
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.4, 10, 0.7, "Problem & Motivation",
             font_size=32, color=WHITE, bold=True)
add_accent_line(slide, 0.8, 1.0, 4)

bullets = [
    "\u2022  Single-image 3D reconstruction: predict a 3D point cloud from one RGB image",
    "\u2022  Models trained on synthetic data (ShapeNet renders) struggle on real-world photos",
    "\u2022  Domain gap: lighting, texture, background, camera differences",
    "\u2022  Goal: bridge this gap using domain adaptation strategies",
    "",
    "\u2022  Key challenge: no ground-truth 3D for real images \u2192 need unsupervised adaptation",
]
add_bullet_slide(slide, 0.8, 1.3, 11.5, 5, bullets, font_size=20, color=LIGHT_GRAY)


# ================================================================
# SLIDE 3: Architecture
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.4, 10, 0.7, "Model Architecture",
             font_size=32, color=WHITE, bold=True)
add_accent_line(slide, 0.8, 1.0, 4)

add_text_box(slide, 0.8, 1.3, 6, 0.5, "Hybrid CNN-Transformer (17.5M parameters)",
             font_size=20, color=ACCENT, bold=True)

arch_bullets = [
    "1.  Encoder: Pretrained ResNet-18 (ImageNet) \u2192 49 image tokens (7\u00d77 grid)",
    "2.  Cross-Attention (2 layers, 8 heads): 2048 learnable query tokens attend to image features",
    "3.  Self-Attention Decoder (4 layers, 8 heads): refine query representations",
    "4.  MLP Head: project each query to (x, y, z) \u2192 2048-point cloud",
    "",
    "\u2022  Differential LR: encoder 0.2\u00d7 base, decoder 5\u00d7 base",
    "\u2022  Cosine LR schedule with 2-epoch warmup",
    "\u2022  Mixed precision training (FP16)",
    "\u2022  Loss: Chamfer Distance (bidirectional)",
]
add_bullet_slide(slide, 0.8, 1.9, 11.5, 5, arch_bullets, font_size=18, color=LIGHT_GRAY)


# ================================================================
# SLIDE 4: Training Setup
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.4, 10, 0.7, "Training Setup",
             font_size=32, color=WHITE, bold=True)
add_accent_line(slide, 0.8, 1.0, 4)

setup_bullets = [
    "\u2022  Dataset: ShapeNet + Cap3D renders (13 categories, 31,832 training samples)",
    "\u2022  Categories: airplane, bench, cabinet, car, chair, display, lamp,",
    "    loudspeaker, rifle, sofa, table, telephone, watercraft",
    "\u2022  Output: 2048-point 3D point cloud per image",
    "\u2022  Batch size: 12  |  Epochs: 40  |  GPU: RTX 4070 Ti Super (16GB)",
    "\u2022  Augmentation: color jitter, gaussian blur, random crop, random erasing",
    "",
    "\u2022  Key bug fixes applied before training:",
    "    \u2013  Flip augmentation: negate point cloud x-coords when image is flipped",
    "    \u2013  Differential learning rates for pretrained encoder vs decoder",
    "    \u2013  ChamferLoss API signature correction",
]
add_bullet_slide(slide, 0.8, 1.3, 11.5, 5.5, setup_bullets, font_size=18, color=LIGHT_GRAY)


# ================================================================
# SLIDE 5: Training Curves
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.4, 10, 0.7, "Training Curves",
             font_size=32, color=WHITE, bold=True)
add_accent_line(slide, 0.8, 1.0, 4)

add_image_safe(slide, f'{VIZ}/training_loss_curve.png', 0.5, 1.3, width=6, height=3.3)
add_image_safe(slide, f'{VIZ}/val_chamfer_distance.png', 6.8, 1.3, width=6, height=3.3)

add_text_box(slide, 0.5, 4.8, 6, 0.5, "Training loss converges smoothly from 0.30 to 0.006",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 6.8, 4.8, 6, 0.5, "Best validation CD: 0.00930 at epoch 40",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 6: Quantitative Results
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.4, 10, 0.7, "Quantitative Results",
             font_size=32, color=WHITE, bold=True)
add_accent_line(slide, 0.8, 1.0, 4)

add_image_safe(slide, f'{VIZ}/results_comparison.png', 1.5, 1.5, width=10, height=3.5)

# Results table as text
table_text = [
    "Metric               Old (1024pts, 100ep)    New (2048pts, 40ep)    Improvement",
    "\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500",
    "CD Mean              0.01522                     0.00980                     -35.6%",
    "F-Score@0.01      0.0196                       0.0424                      +116%",
    "F-Score@0.02      0.1211                        0.1958                      +61.7%",
    "F-Score@0.05      0.5865                       0.6615                      +12.8%",
]
add_bullet_slide(slide, 1.0, 5.2, 11, 2, table_text, font_size=13, color=LIGHT_GRAY,
                 font_name='Consolas')


# ================================================================
# SLIDE 7: Synthetic Results — Best
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.4, 10, 0.7, "Synthetic Test Results \u2014 Best Samples",
             font_size=32, color=WHITE, bold=True)
add_accent_line(slide, 0.8, 1.0, 4)

add_image_safe(slide, f'{VIZ}/synth_airplane.png', 0.3, 1.3, width=6.2, height=1.7)
add_image_safe(slide, f'{VIZ}/synth_car1.png', 6.8, 1.3, width=6.2, height=1.7)
add_image_safe(slide, f'{VIZ}/synth_rifle1.png', 0.3, 3.2, width=6.2, height=1.7)
add_image_safe(slide, f'{VIZ}/synth_lamp.png', 6.8, 3.2, width=6.2, height=1.7)

add_text_box(slide, 0.8, 5.1, 11.5, 0.8,
             "Model captures overall shape well for compact objects (airplanes, cars, rifles, lamps). "
             "Predicted point clouds (blue) closely match ground truth (green).",
             font_size=15, color=LIGHT_GRAY)


# ================================================================
# SLIDE 8: Synthetic Results — More
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.4, 10, 0.7, "Synthetic Test Results \u2014 More Samples",
             font_size=32, color=WHITE, bold=True)
add_accent_line(slide, 0.8, 1.0, 4)

add_image_safe(slide, f'{VIZ}/synth_cabinet.png', 0.3, 1.3, width=6.2, height=1.7)
add_image_safe(slide, f'{VIZ}/synth_car2.png', 6.8, 1.3, width=6.2, height=1.7)
add_image_safe(slide, f'{VIZ}/synth_rifle2.png', 0.3, 3.2, width=6.2, height=1.7)

add_text_box(slide, 6.8, 3.2, 6, 1.5,
             "Consistent performance across categories:\n"
             "\u2022  Cabinets: dense box-like shapes\n"
             "\u2022  Cars: elongated body captured\n"
             "\u2022  Rifles: thin elongated structure",
             font_size=16, color=LIGHT_GRAY)


# ================================================================
# SLIDE 9: Domain Adaptation — TTA
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.4, 10, 0.7, "Domain Adaptation: Test-Time Augmentation",
             font_size=32, color=WHITE, bold=True)
add_accent_line(slide, 0.8, 1.0, 4)

tta_bullets = [
    "\u2022  Strategy: apply 10 random augmentations at inference, average predictions",
    "\u2022  Augmentations: color jitter, random crop, horizontal flip (with x-coord correction)",
    "\u2022  No retraining needed \u2014 applied at test time only",
    "\u2022  Reduces prediction variance (spread), producing tighter point clouds",
]
add_bullet_slide(slide, 0.8, 1.3, 11.5, 2.5, tta_bullets, font_size=18, color=LIGHT_GRAY)

add_image_safe(slide, f'{VIZ}/real_airplane.png', 0.3, 3.8, width=6.2, height=1.7)
add_image_safe(slide, f'{VIZ}/real_rifle1.png', 6.8, 3.8, width=6.2, height=1.7)

add_text_box(slide, 0.8, 5.7, 11.5, 0.5,
             "Blue = No TTA (scattered)  |  Orange = With TTA (tighter, more coherent)",
             font_size=14, color=GOLD, alignment=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 10: Real-World Results
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.4, 10, 0.7, "Real-World Inference Results",
             font_size=32, color=WHITE, bold=True)
add_accent_line(slide, 0.8, 1.0, 4)

add_image_safe(slide, f'{VIZ}/real_car.png', 0.3, 1.3, width=6.2, height=1.7)
add_image_safe(slide, f'{VIZ}/real_monitor.png', 6.8, 1.3, width=6.2, height=1.7)
add_image_safe(slide, f'{VIZ}/real_lamp.png', 0.3, 3.2, width=6.2, height=1.7)
add_image_safe(slide, f'{VIZ}/real_rifle2.png', 6.8, 3.2, width=6.2, height=1.7)

add_text_box(slide, 0.8, 5.1, 11.5, 0.8,
             "Model generalizes from synthetic training to real photographs. "
             "Background removal (rembg) + TTA significantly improves reconstruction quality.",
             font_size=15, color=LIGHT_GRAY)


# ================================================================
# SLIDE 11: Qualitative Comparison — Old vs New on Chairs
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.4, 10, 0.7, "Qualitative Comparison: 1024-pt vs 2048-pt",
             font_size=32, color=WHITE, bold=True)
add_accent_line(slide, 0.8, 1.0, 4)

add_text_box(slide, 0.3, 1.2, 6, 0.4, "Old Model (1024 pts, 100 epochs)",
             font_size=18, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 6.8, 1.2, 6, 0.4, "New Model (2048 pts, 40 epochs)",
             font_size=18, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)

add_image_safe(slide, f'{VIZ}/old_chair1.png', 0.3, 1.7, width=6.2, height=1.6)
add_image_safe(slide, f'{VIZ}/limitation_chair1.png', 6.8, 1.7, width=6.2, height=1.6)
add_image_safe(slide, f'{VIZ}/old_chair2.png', 0.3, 3.5, width=6.2, height=1.6)
add_image_safe(slide, f'{VIZ}/limitation_chair2.png', 6.8, 3.5, width=6.2, height=1.6)

add_text_box(slide, 0.8, 5.3, 11.5, 1.0,
             "Old model captures chair structure better (100 epochs of refinement). "
             "New model produces denser but more scattered predictions on thin structures. "
             "With more training, 2048-pt model would likely surpass the old model on chairs too.",
             font_size=15, color=LIGHT_GRAY)


# ================================================================
# SLIDE 12: Limitations & Future Work
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.4, 10, 0.7, "Limitations & Future Work",
             font_size=32, color=WHITE, bold=True)
add_accent_line(slide, 0.8, 1.0, 4)

limit_bullets = [
    "Limitations:",
    "  \u2022  Thin structures (chair legs, lamp arms) remain challenging",
    "  \u2022  Chamfer Distance allows spread-out solutions \u2014 doesn\u2019t penalize off-surface points",
    "  \u2022  Single-view ambiguity: occluded parts are guessed, not seen",
    "  \u2022  2048-pt model only trained 40 epochs (vs 100 for 1024-pt baseline)",
    "",
    "Future Work:",
    "  \u2022  Train 2048-pt model for 100+ epochs for fair comparison",
    "  \u2022  Add Earth Mover\u2019s Distance (EMD) or surface-aware loss for tighter reconstructions",
    "  \u2022  Multi-view fusion: combine predictions from multiple viewpoints",
    "  \u2022  AdaIN style transfer for stronger domain adaptation",
    "  \u2022  Mesh reconstruction from predicted point clouds",
]
add_bullet_slide(slide, 0.8, 1.3, 11.5, 5.5, limit_bullets, font_size=18, color=LIGHT_GRAY)


# ================================================================
# SLIDE 13: Conclusion
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.4, 10, 0.7, "Conclusion",
             font_size=32, color=WHITE, bold=True)
add_accent_line(slide, 0.8, 1.0, 4)

concl_bullets = [
    "\u2022  Built a hybrid CNN-Transformer model for single-image 3D point cloud reconstruction",
    "",
    "\u2022  Trained on 13 ShapeNet categories (31,832 samples) with 2048-point output",
    "",
    "\u2022  Achieved 35.6% lower Chamfer Distance and 2\u00d7 better F-Score@0.01 vs baseline",
    "",
    "\u2022  Test-Time Augmentation successfully bridges synthetic-to-real domain gap",
    "    without any real-world 3D supervision",
    "",
    "\u2022  Model generalizes to real photographs across multiple object categories",
    "",
    "\u2022  Identified key training pitfalls: flip augmentation alignment, differential LR,",
    "    and loss function API verification are critical for point cloud reconstruction",
]
add_bullet_slide(slide, 0.8, 1.3, 11.5, 5.5, concl_bullets, font_size=19, color=LIGHT_GRAY)


# ================================================================
# SLIDE 14: Thank You
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 1.5, 2.5, 10, 1.0, "Thank You!",
             font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, 4, 3.7, 5.3)

add_text_box(slide, 1.5, 4.0, 10, 0.6, "Questions?",
             font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1.5, 5.2, 10, 0.5, "Mrinal Bharadwaj  |  AI 535  |  March 2026",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# Save
output_path = './AI535_3D_Reconstruction_Presentation.pptx'
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
